from __future__ import annotations

from pathlib import Path

from convert_to_md.context import ConvertContext
from convert_to_md.converters.base import BaseConverter
from convert_to_md.ir import DocumentIR, ImageAsset


class PptxConverter(BaseConverter):
    name = "pptx"
    kinds = ("pptx",)

    def convert(self, path: Path, ctx: ConvertContext) -> DocumentIR:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(path))
        doc = DocumentIR(source=path, title=path.stem)
        image_i = 0

        for idx, slide in enumerate(prs.slides, start=1):
            title = _slide_title(slide) or f"Slide {idx}"
            doc.heading(title, level=1)

            # collect text / tables / images in shape order
            for shape in slide.shapes:
                if shape.has_table:
                    rows = []
                    table = shape.table
                    for r in table.rows:
                        rows.append([c.text_frame.text.replace("\n", " ").strip() for c in r.cells])
                    if len(rows) > ctx.max_table_rows:
                        rows = rows[: ctx.max_table_rows]
                    doc.table(rows)
                    continue

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_i += 1
                        ext = f".{image.ext}" if image.ext else ".png"
                        asset = ImageAsset(
                            data=image.blob,
                            filename=f"slide{idx}_image_{image_i}{ext}",
                            alt=getattr(shape, "name", "") or f"image_{image_i}",
                            content_type=image.content_type,
                        )
                        doc.image(asset)
                    except Exception:
                        pass
                    continue

                if shape.has_text_frame:
                    # skip title shape text if already used as heading
                    if _is_title_shape(shape, title):
                        continue
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if not text:
                            text = (para.text or "").strip()
                        if text:
                            doc.paragraph(text)

            if idx < len(prs.slides):
                doc.hr()

        return doc


def _slide_title(slide) -> str | None:
    if slide.shapes.title is not None:
        t = slide.shapes.title.text
        if t and t.strip():
            return t.strip()
    return None


def _is_title_shape(shape, title: str) -> bool:
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    return bool(title) and text == title
